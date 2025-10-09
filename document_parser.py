import fitz  # PyMuPDF
import pandas as pd
import docx
from pptx import Presentation

def parse_document(filepath):
    ext = filepath.split('.')[-1].lower()
    if ext == "pdf":
        doc = fitz.open(filepath)
        return [page.get_text() for page in doc]
    elif ext == "csv":
        df = pd.read_csv(filepath)
        return [df.to_string()]
    elif ext == "docx":
        doc = docx.Document(filepath)
        return ["\n".join([para.text for para in doc.paragraphs])]
    elif ext == "pptx":
        prs = Presentation(filepath)
        return [" ".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])]
    elif ext in ["txt", "md"]:
        with open(filepath, "r", encoding="utf-8") as f:
            return [f.read()]
    return []